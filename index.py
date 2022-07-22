#!/usr/bin/env python3
# A REST API to extract data from PowerPoint presentations
# Robert Ryder, 31 May 2022
from flask import * # not included; pip3 install flask
import gif_compositor
import pptx # not included; pip3 install python-pptx
import lxml.etree # not included, but a dependency of pptx
import dotenv # not included; pip3 install python-dotenv
import zipfile
import shutil
import json
import time
import hashlib
import os
import os.path
import sys
import threading
import queue
import logging

"""
Some information on the new XML-based PPTX format that PowerPoint uses
I. Distances
Distances are given in DXA, which is 1/20 (0.05) pt. Thus, 1pt = 20dxa, 1 inch = 72pt = 1440dxa, and 1cm = 0.3937007874in = 566.929133856dxa
Font sizes appear to be given in half-points to avoid non-integer values


II. File structure and data format as pertaining to this script
All PPTX files are just ZIP files with a few extra features. Any archive utility should be able to extract them, but files re-zipped
with a normal archive utility may or may not work again in PowerPoint. Extracting the file, two items will be produced: a file named "[Content Types].xml",
and a directory named "ppt". The latter will hold the contents of the file.

All video, images, and audio will be stored in the directory ppt/media; audio and video files will be named  mediaN.[extension] and images imageN.[extension],
 where N is the order in which the file was added to the presentation. 

Each slide will have a file stored in /ppt/slides/_rels/ named "slideN.xml.rels" containing all external references in the "Target" attribute of <Relationship> and/or <Relationships> tags.


III. Software requirements for this program
1. Python 3.8+
2. The following Python modules that are not included with the standard distribution: flask, python-pptx, lxml (the Python bindings to the C libraries lxml2 and libxslt.
All of these are readily installable using PIP


IV. What this program does

1. The /form endpoint (POST, multipart/form-data, file must be in key named "pres", API authentication key must be in key named "key")
This returns a ZIP archive containing all embedded media from the attached PowerPoint file. Also in this archive will be a file named webpptx-metadata.json, 
which will list the aspect ratio, the speaker notes, and all video and audio files embedded in the presentation. Each index of the the array value for each key corresponds 
to a slide; multiple values per slide are stored in sub-arrays. Embedded images will be named imageN.[extension] (N denotes that this was the Nth image added to the presentation),
and audio or video will be named mediaN.[extension].

If the request does not contain values for the API key ("key") or the presentation file ("pres"), a 400 bad request error will result. If the attached file cannot be read or is 
not a .pptx file, a 422 unprocessable entity error will result. The API key included in the request must be valid, otherwise, the server will respond with a 403 forbidden error.

2. The /animate endpoint (POST, mutlpart/form-data, presentation must be in key named "pres", various other files)
If given a PPTX presentation and the images (from ConvertAPI) of each slide, this will extract all the animated GIF files from the presentation and overlay them on each slide. It returns a ZIP archive
containing animated GIFs of each slide, the animations drawn precisely where their still-image versions were. If there are no animations on a particular slide, the original PNG image will be returned instead.

2. Environment variables
Environment variables can either be specified in the standard way or in a .env file.

a. API_KEY: This is the string that all requests must contain to be processed. It must be set, or the program will exit

b. TEMP_DIR: This can be used to control where temp files are placed. If not set, the current working directory will be used.


V. TODOs/Concerns/Notes/etc.
1. To extend the functionality beyond what python-pptx offers, I suspect this code does some of the same things that that library does under the hood, which is redundant.

2. The presentation must be included with every request, which consumes a lot of bandwidth. The PowerPoint file I used to put this through its paces was 105 MB.
For two endpoints, that's 210 MB, plus the ZIP archive of media sent back the other way.

3. The logging is kind of ad-hoc.

4. If you stop this program, any temp files scheduled for cleanup will not be deleted- you will have to delete them manually.
"""

dotenv.load_dotenv()  # initialize environment variables specific to this program

# This could be a configuration file, but the information stays constant enough that I don't think it's necessary.

# Only the new XML-based PowerPoint files are supported. Initially, I thought that python-pptx would support the specialty PowerPoint formats, such as .ppsx/PowerPoint Show (hence why this is a list),
# but it does not appear to feature such support
ALLOWED_EXTENSIONS = ['.pptx']

ALLOWED_IMAGES = ['.png', '.gif', '.jpg', '.tif', '.tiff']  # More formats could be supported, but anyone that uses OpenEXR or TGA or PXR or the like can deal with that themselves. With such obscure formats, they will probably be used to it.

# Set where to put temp files. Theoretically, this can be anywhere, but . , a subdirectory of . , and /tmp make the most sense.
# display a more obvious error message than 'KeyError: os.environ has no attribute...'
try:
    TEMP_DIR = os.environ['TEMP_DIR']
except KeyError:
    sys.stderr.write('Temporary file location is not defined')  # flushed anyway on exit()
    exit(-1)

# Used to verify authenticity with upstream API

try:
    API_KEY = os.environ['API_KEY']  # All-caps with underscores, in the tradition of Unix environment variables
except KeyError:
    sys.stderr.write('Unable to obtain secure key for API')
    exit(-1)

# Initialize logging
logging.basicConfig(format='%(asctime)s %(levelname)s %(filename)s %(funcName)s:%(lineno)d %(name)s %(message)s')
log = logging.getLogger(__name__)
log.setLevel(logging.DEBUG)
log.info('test')

# Look for URLs like these in the .xml.rels files for each slide.
# It is very important that you add a trailing slash (/) to each URL template to prevent something like https://www.youtube.com.otherdomain.malware.xyz/download.txt.vbs.php
# Sourced from the following: https://support.microsoft.com/en-us/office/video-and-audio-file-formats-supported-in-powerpoint-d8b12450-26db-4c7b-a5c1-593d3418fb59

WEB_VIDEO_URLS = ['.mov',  # Apple QuickTime
                  '.qt',
                  '.mp4',  # MPEG 1/2/3/4              
                  '.m4v',
                  '.mpg',
                  '.mpeg',
                  '.mpe',
                  '.m15',
                  '.m75',
                  '.m2v',
                  '.ts',
                  '.wmv',  #Windows Media Video
                  '.dvi',  # DVI video
                  '.avi',  # AVI video
                  '.vfw',
                  '.asf',
                  'https://www.youtube.com/',
                  'https://player.vimeo.com/',
                  'https://dailymotion.com/']

EMBEDDED_AUDIO_TYPES = [ '.aif', # Apple Lossless
                         '.aiff',
                         '.au', # Sun Microsystems audio format
                         '.snd',
                         '.mid',  # MIDI
                         '.midi',
                         '.mp3',  # MPEG
                         '.mpga',
                         '.m4a',  # M4A is both MPEG and an Apple format
                         '.wav', 
                         '.wave',  # WAVE
                         '.bwf',  #'.BWF',
                         '.aa', # Amazon Audible
                         '.aax',
                         '.wma',
                         '.aac',  # Various Apple formats
                         '.caf',
                         '.m4r',
                         '.ac3',
                         '.eac3']


SEARCH_CRAWLER_TOKENS = [
    'APIs-Google',
    'AdsBot-Google',
    'Mediapartners-Google',
    'Googlebot',
    'FeedFetcher-Google',
    'Google-Read-Aloud',
    'DuplexWeb-Google',
    'googleweblight',  # Thank you, Google, for such a fragmented ecosystem of search crawlers, including many that don't respect robots.txt entreaties that they leave your data alone. I hope this is all of them.
    'bingbot',
    'duckduckbot'
]  # Don't let search engines index the API


# Initially, this used one delete queue that was only ever called from the /form endpoint. However, to ensure files are deleted more promptly
# but that no signals are crossed, two queues are used. In a multithreaded program, these should perhaps be changed to queue.Queue() objects.
to_delete_form = []  # store temp files scheduled for deletion for the form endpoint
to_delete_animate = []  # store temp files scheduled for deletion for the animate endpoint

application = Flask(__name__)
application.secret_key = os.urandom(16)


# Checks whether a file can be uploaded, whether it has the correct extension for a presentation
def allowable_file(filename): return os.path.splitext(filename)[-1] in ALLOWED_EXTENSIONS

#Does the same as the above, but for images
def allowable_image(filename): return os.path.splitext(filename)[-1] in ALLOWED_IMAGES


# This could potentially be refactored into its own module or combined with gif_compositor to form a "utils" file
class AsynchronousProcessingThread(threading.Thread):
    # This is used to run the /animate endpoint in the background, so that the request can periodically send back bytes to the effect of 'continue receiving data'
    # Heroku routers terminate the connection after thirty seconds of silence, this is used to ensure that that timeout is never reached. I used a thread rather than
    # async/await because the response.zip file is written to disk- there is nothing that this returns which I can await. Even if there was, I couldn't use async/await, because
    # await (true to its name) *waits*. The only reason this is done asynchronously in the first place is that I can't wait for the full response.
    # This didn't work either: it merely changed my H12 timeout errors into H18 server connection interrupted errors: the Heroku router cut the connection anyway.
    def initialize_queue(self):
        # Set up queues and log a message to ensure that the thread started properly.
        log.debug('worker thread initialized')
        self.job_queue = queue.Queue()
        self.available_results_animation = queue.Queue()
        self.available_results_form = queue.Queue()
        
    def enqueue_job(self, job_type, job_id, hashpath, hashpath_image=None, image_files=None, presentation=None):
        #Create objects used in run()- I need objects from the request to be accessible there, and run() is only called from within the threading module;
        # therefore, store data within the class using this method
        assert (hashpath_image and image_files) or presentation  # one of these must be set
        if job_type == 'A':
            log.debug('enqueued animation job')
            self.job_queue.put(dict(job_type=job_type, job_id=job_id, hashpath=hashpath, hashpath_image=hashpath_image, image_files=image_files))
        elif job_type == 'N':
            log.debug('enqueued notes/media job')
            self.job_queue.put(dict(job_type=job_type, job_id=job_id, hashpath=hashpath, presentation=presentation))
            
    def execute_notes_media_job(self):
        log.debug('starting notes/media job')
        p = self.presentation
        response_path = TEMP_DIR + self.job_id + 'notesmedia'  # + '.zip'
        pptx_data = {}
        notes = []
        videos = []
        tmp = []
        audio = []
        tmp_audio = []
        hashpath = self.hashpath
        # Get speaker notes for each slide
        for i in p.slides:
            if i.has_notes_slide:
                notes.append(i.notes_slide.notes_text_frame.text)
            else:
                notes.append('')  # keep the array of notes aligned to a zero-based index of slide numbers; add empty string if no notes exist

        # New PowerPoint files are just .ZIP archives; extract them to obtain needed data. Notably, however, files extracted and
        # re-zipped often will not work, but the special nature only works in that direction.
        with zipfile.ZipFile(hashpath) as z:
            z.extractall(hashpath + '_extracted')
            
        rels_dir = hashpath + '_extracted/ppt/slides/_rels/'
        for i in sorted(os.listdir(rels_dir)):  # use sorted() to ensure notes are in the proper order. Per the documentation's philosophical-sounding language for os.listdir(): "if order matters, you must impose it."
            if 'slide' in i and '.xml.rels' in i:  # exclude files other than slideN.xml.rels
                # I tried this with "Relationship," "Relationships," etc. to try and isolate what I wanted; none of it worked properly.
                # I think it's impossible to evade iterating over everything. However, it's a small and simple file.
                elements = lxml.etree.parse(rels_dir + i).xpath('/*')
                for j in elements:
                    for k in j.getchildren():
                        for l in WEB_VIDEO_URLS:
                            if l in k.attrib['Target'].lower() and k.attrib['Target'] not in tmp:  # 'Target' attribute looks like a video URL. Prevent duplicates.
                                tmp.append(k.attrib['Target'])  # add to videos

                        for l in EMBEDDED_AUDIO_TYPES:
                            if l in k.attrib['Target'].lower() and k.attrib['Target'] not in tmp_audio:  # 'Target' attribute looks like a reference to an audio file
                                tmp_audio.append(k.attrib['Target'])  # add to audio

            videos.append(tmp)  # use array tmp to put all videos from each slide into a sub-array for organization purposes
            tmp = []
            audio.append(tmp_audio)
            tmp_audio = []

        # Include aspect ratio - this matters for front-end corner cases. I am fairly certain 5:4 presentations do not exist.
        if (p.slide_width / float(p.slide_height)) < 1.4: # 4:3 is 1.333...
            pptx_data['aspect_ratio'] = '4:3'
        else:  # 16:9 is 1.777...
            pptx_data['aspect_ratio'] = '16:9'

        # Create orderly JSON data for response
        pptx_data['notes'] = notes
        pptx_data['videos'] = videos
        pptx_data['audio'] = audio
        # clean up extracted files
        del p  # Garbage collection should really do this instead once the variable is out of scope, but having this here seems to improve memory usage.

        # Extract Media
        # Create unique directory of contents for Office/PowerPoint XML file .zip
        #zipfile.ZipFile(hashname).extractall(hashname + '_extracted')
        #response_name = hashlib.md5(str(time.time()).encode('ascii')).hexdigest() + '.zip'  # create a zip file with a distinct name for the response.
        media_path = hashpath + '_extracted/ppt/media/'
        response_path = TEMP_DIR + self.job_id + 'notesmedia'
        # Files are copied to put them in their own folder at the top level of TEMP_DIR. The reason that is important is deletion-queue logistics- the media path cannot
        # be scheduled for deletion until the results are fetched, which happens outside of this function. At that point, only the path to the media folder is available; the
        # outer folders of [hash].pptx_extracted/ppt will not be included. That, in turn, means that they won't be scheduled for deletion, because they won't be added to the queue
        # along with the media queue. This is a less elegant solution than I would like, but every alternative is worse. I could, for example, pass the other path down to the request handler,
        # but that would be far less streamlined.
        shutil.copytree(media_path, response_path)
        with open(response_path + '/webpptx-metadata.json', 'x') as f:
            json.dump(pptx_data, f, indent=2)
        #shutil.copytree(hashpath + '_extracted/ppt/media/', response_path)
        log.debug('media path is: ' + media_path)
        shutil.rmtree(hashpath + '_extracted')
        self.available_results_form.put(dict(job_id=self.job_id, content_path=response_path))
        
    def execute_animation_job(self):
        log.debug('starting animation job')
        p = pptx.Presentation(self.hashpath)
        SLIDES_DIRECTORY = TEMP_DIR + self.job_id  # 'response-{}'.format(hashlib.md5(str(time.time()).encode('ascii')).hexdigest())
        os.mkdir(SLIDES_DIRECTORY)  # used to store slides
        # Track index in the foreach loop. Perhaps the right way to do this is with a normal for accessing the array of objects,
        # but having the object controlled by the loop greatly improves readability and conciseness.
        slide_number = 0  
        for i in p.slides:
            gifs = []
            frame_counter = -1  # Used to determine the number of frames in the composite slide.
            idx = 0
            for j in i.shapes:
                if j.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE and j.image.ext == 'gif':
                    # I want to do this using fromBytes, but that requires that I specify the size.
                    # I don't know the size without reading .size on the object I can't create unless I know the size.
                    # I tried using BytesIO, but then the images all looked like they had TV overscan; the margins were cut off.
                    gif = gif_compositor.GIFContainer(j.image.blob, self.hashpath_image + '-{}-'.format(idx) + '_gif.gif', j.left, j.top, j.width, j.height)

                    image = gif_compositor.PIL.Image.open(self.hashpath_image)  # written to a file at hashpath_image in the GIFContainer constructor
                    image_width, image_height = image.size

                    # Convert the dimensions of the image that PowerPoint gives in DXA to pixels, scaling to the fit the size of the attached image
                    gif.scaleXPosition(image_width / p.slide_width)
                    gif.scaleYPosition(image_height / p.slide_height)
                    gif.scaleWidth(image_width / p.slide_width)
                    gif.scaleHeight(image_height / p.slide_height)

                    print(gif.img.info['duration'])
                    if gif.img.n_frames > frame_counter:
                        frame_counter = gif.img.n_frames  # set this to the maximum number of frames in any GIF on this slide

                    gifs.append(gif)
                    idx += 1

            try:
                gifs = sorted(gifs, key=lambda g: g.framedelay)  # sort by framedelay. 
                if len(gifs) > 0 and frame_counter > 1:  # Ensure GIFs exist and that they have multiple frames of animation
                    gif_compositor.compose_gifs(self.image_files[slide_number], gifs, '{}/slide{}.gif'.format(SLIDES_DIRECTORY, slide_number+1), frame_counter)  # Generate the composite GIF of the slide from its constituent parts.

                else: 
                    shutil.copyfile(self.image_files[slide_number], '{}/slide{}.png'.format(SLIDES_DIRECTORY, slide_number+1))  # Copy original, non-animated slide into response.

            except IndexError:
                # we didn't get enough images to represent the slides
                abort(422)  # client's error, use 4XX
            slide_number += 1

        # Do this here, synchronous with the above
        for i in self.image_files: os.remove(i)

        to_delete_animate.append(self.hashpath)
        self.available_results_animation.put(dict(job_id=self.job_id, content_path=SLIDES_DIRECTORY))
        log.debug('Done inside animation worker thread')

    def run(self):
        # In the initial version of this program that did this asynchronously, run() was execute_single_job(). For that reason, these parameters are passed
        # through the class rather than as keyword arguments.
        print('starting thread')
        while True:
            job = self.job_queue.get()
            self.job_id = job['job_id']
            self.hashpath = job['hashpath']
            self.RESPONSE_PATH = TEMP_DIR + hashlib.md5(str(time.time()).encode('ascii')).hexdigest()
            if job['job_type'] == 'A':  # A is animation, N is notes/media
                self.hashpath_image = job['hashpath_image']
                self.image_files = job['image_files']
                self.execute_animation_job()
            elif job['job_type'] == 'N':
                self.presentation = job['presentation']
                self.execute_notes_media_job()
                
            del job


# Initialize worker thread
worker = AsynchronousProcessingThread()
worker.initialize_queue()
worker.start()


@application.route('/')
def confirm_activity():
    # Check that this isn't a search engine. Robots.txt is voluntary, but this is enforced
    agent = request.headers.get('User-Agent')
    for i in SEARCH_CRAWLER_TOKENS:
        if i in agent: abort(404)  # we don't exist
    
    return '<h1>This application is working properly</h1>'  # sanity-check deployment


@application.route('/upload', methods=['POST'])
def animate_presentation():
    try:
        if request.form['key'] != API_KEY: abort(403)
    except KeyError:
        abort(403)

    # Check whether a file has the proper extension and that it exists
    fileitem = request.files['pres']
    fname = fileitem.filename
    if fname == '':
        sys.stderr.write('no presentation file')
        sys.stderr.flush()
        
    # Use a hash of the time to ensure a unique filename and thus prevent duplicate files or cross-signals
    # Use MD5 because this hash isn't used for anything especially security-related; you don't need password-grade cryptography.
    # MD5 should be simpler and thus faster than other hash algorithms.
    if fileitem and allowable_file(fname):
        hashname = hashlib.md5(str(time.time()).encode('ascii')).hexdigest() + '.pptx'
        fileitem.save(TEMP_DIR + hashname)
    else:
        abort(422)

    hashpath = TEMP_DIR + hashname  # location of PPTX file

    request_data = dict(request.files)  # create a mutable dict() from the immutable listing of attached files
    request_data.pop('pres')
    log.debug('req data: ' + str(request_data))

    if len(request_data) == 0:
        # No files attached. This is the client's fault, so use 4XX error code.
        abort(400)  # no entity, not a 422
        
    image_files = []  # store the locations of temporary image files

    # Download attached files, which should be named slideN.png starting at 0 or 1- it doesn't matter.
    for file_key_name in sorted(list(request_data.keys())):
        fileitem = request.files[file_key_name]
        fname = fileitem.filename
        if fname == '':
            sys.stderr.write('no presentation file')
            sys.stderr.flush()
    
        if fileitem and allowable_image(fname):
            hashpath_image = TEMP_DIR + 'image-' + hashlib.md5(str(time.time()).encode('ascii')).hexdigest() + '.png'  # random name for the files
            image_files.append(hashpath_image)
            fileitem.save(hashpath_image)

    # If this fails, it is a client-side error, so use 4xx rather than 5xx error codes.
    try:
        p = pptx.Presentation(hashpath)
    except pptx.exc.PackageNotFoundError:
        # not a PowerPoint presentation
        abort(422)
    except ValueError:
        # If the file  is a specialty PowerPoint file (e. g., ppsx), a ValueError will result
        # The extension would have to be mistaken and the various content-type metadata incongruous.
        abort(422)

    job_id = 'job-' + hashlib.md5(str(time.time()).encode('ascii')).hexdigest()  # Used to match output from /available_jobs with its initial requests.
    
    # Everything above this line is required to ensure that the request is valid and that the worker thread will succeed.
    # Having confirmed that everything is in the proper order, we can punt to the asynchronous worker thread here
    
    worker.enqueue_job('A', job_id, hashpath, hashpath_image, image_files)  # Do the animation thread first- most presentations don't contain GIFs, and animation will therefore usually finish more quickly.
    worker.enqueue_job('N', job_id, hashpath, presentation=p)  # enqueue notes job.
    
    return json.dumps(dict(jobID=job_id))  # match results to initial requests later


@application.route('/form-results')
def return_form_results():
    # Do it this way rather than using a foreach and initializing the array just in case something is added while this executes.
    # Perhaps this should use locking, but this works and I don't want to risk any exceptions on other routes that can't access this.
    # It is very important that this produce net-zero temp files. This happens before authentication because it doesn't reveal any data- it's just general maintenance.
    log.info('sending available form results')

    agent = request.headers.get('User-Agent')
    for i in SEARCH_CRAWLER_TOKENS:
        if i in agent: abort(404)  # we don't exist
        

    for i in range(len(to_delete_form)):
        try:
            if os.path.exists(to_delete_form[i]):
                if os.path.isdir(to_delete_form[i]):
                    shutil.rmtree(to_delete_form[i])
                else:
                    os.remove(to_delete_form[i])

                del to_delete_form[i]
                
        # If the file doesn't exist, all to the good- not a problem
        except IndexError:
            log.warn('IndexError in deleting temporary file')
        except FileNotFoundError:
            log.warn('FileNotFoundError in cleaning temporary file')

    try:
        if request.form['key'] != API_KEY: abort(403)
    except KeyError:
        abort(403)

    if worker.available_results_form.qsize() == 0: return ''  # Nothing to show. I was, and to certain extent remain, conflicted about whether to use 410 Gone here.

    response_path = TEMP_DIR + 'response-' + str(int(time.time())) + '.zip'

    # create final response ZIP containing subfolders of each job. Job ID and content path are passed separately because the job ID is determined when the request is processed (cannot be done here)
    # and content_path cannot encode it without including TEMP_DIR and other data that will be useless to the client.
    with zipfile.ZipFile(response_path, mode='x') as response:
        while worker.available_results_form.qsize() > 0:  # get everything from the results queue
            item = worker.available_results_form.get()  # create a copy here because get() removes the item from the queue and it would be rendered inaccessible
            for i in os.listdir(item['content_path']):
                response.write(item['content_path'] + '/' + i, arcname='response/' + item['job_id'] + '/' + i)
                to_delete_form.append(item['content_path'])  # schedule this for deletion

        response.close()  # write headers and other metadata to ZIP file.
            
    to_delete_form.append(response_path)
    return send_file(response_path)  # send response


@application.route('/animation-results')
def return_animation_results():
    log.info('starting available jobs route')
    agent = request.headers.get('User-Agent')
    for i in SEARCH_CRAWLER_TOKENS:
        if i in agent: abort(404)  # we don't exist
        
    for i in range(len(to_delete_animate)):
        try:
            if os.path.exists(to_delete_animate[i]):
                if os.path.isdir(to_delete_animate[i]):
                    shutil.rmtree(to_delete_animate[i])
                else:
                    os.remove(to_delete_animate[i])

                del to_delete_animate[i]
        except IndexError:
            log.warn('IndexError in deleting temporary file')
        except FileNotFoundError:
            log.warn('FileNotFoundError in deleting temporary file')

    # authentication
    try:
        if request.form['key'] != API_KEY: abort(403)
    except KeyError:
        abort(403)
        
    if worker.available_results_animation.qsize() == 0:
      return ''  # nothing available to see here; send an empty response.

    response_path = TEMP_DIR + 'response-' + str(int(time.time())) + '.zip'
    with zipfile.ZipFile(response_path, mode='x') as f:
        while worker.available_results_animation.qsize() > 0:
            item = worker.available_results_animation.get()
            to_delete_animate.append(item['content_path'])  # not deleted until after the response is sent
            print('item path is: ', item['content_path'])
            for i in os.listdir(item['content_path']):
                f.write(item['content_path'] + '/' + i, 'response/' + item['job_id'] + '/' + i)  # use arcname to avoid revealing server temp file structure.
                
        f.close()

    to_delete_animate.append(response_path)
    
    return send_file(response_path)


if __name__ == '__main__':
    # This will not run in production, as it will probably be imported as a module there
    application.run()
